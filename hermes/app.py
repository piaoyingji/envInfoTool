from __future__ import annotations

import csv
import io
import os
import re
import subprocess
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

from fastapi import FastAPI
from minio import Minio
from pydantic import BaseModel


app = FastAPI(title="OneCRM Hermes Agent", version="0.1.0")


class SourceFile(BaseModel):
    id: str = ""
    filename: str = ""
    objectKey: str
    sha256: str = ""
    contentType: str = ""
    sizeBytes: int = 0
    relativePath: str = ""
    clientModifiedAt: str = ""
    uploadedAt: str = ""


class VpnIngestRequest(BaseModel):
    bucket: str
    files: list[SourceFile]


def env(name: str, default: str = "") -> str:
    return os.environ.get(name, default)


def minio_client() -> Minio:
    secure = env("ONECRM_MINIO_SECURE", "false").lower() in {"1", "true", "yes", "on"}
    return Minio(
        env("ONECRM_MINIO_ENDPOINT", "onecrm-minio:9000"),
        access_key=env("ONECRM_MINIO_ACCESS_KEY", "onecrm"),
        secret_key=env("ONECRM_MINIO_SECRET_KEY", "onecrm_minio_pass"),
        secure=secure,
    )


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.post("/internal/vpn-ingest")
def ingest_vpn_sources(payload: VpnIngestRequest) -> dict[str, Any]:
    client = minio_client()
    fragments: list[dict[str, Any]] = []
    warnings: list[str] = []
    for source in payload.files:
        filename = source.filename or Path(source.objectKey).name
        if filename.lower().endswith(".dmp"):
            warnings.append(f"Rejected dmp file: {filename}")
            continue
        try:
            response = client.get_object(payload.bucket, source.objectKey)
            data = response.read()
            response.close()
            response.release_conn()
            fragments.extend(parse_source(filename, source, data, warnings))
        except Exception as exc:
            warnings.append(f"{filename}: {exc}")
    source_meta = summarize_source_meta(fragments)
    precedence_summary = compose_precedence_summary(source_meta)
    raw_text = compose_raw_text(fragments, warnings)
    return {
        "rawText": raw_text,
        "sourceRawText": raw_text,
        "fragments": fragments,
        "sourceMeta": source_meta,
        "precedenceSummary": precedence_summary,
        "warnings": warnings,
    }


def parse_source(filename: str, source: SourceFile, data: bytes, warnings: list[str]) -> list[dict[str, Any]]:
    suffix = Path(filename).suffix.lower()
    if suffix == ".dmp":
        warnings.append(f"Rejected dmp file: {filename}")
        return []
    if suffix == ".zip":
        return parse_zip(filename, source, data, warnings)
    if suffix == ".pdf":
        return [fragment(source, "pdf", extract_pdf_text(data))]
    if suffix in {".docx"}:
        return [fragment(source, "docx", extract_docx_text(data))]
    if suffix in {".xlsx", ".xlsm"}:
        return [fragment(source, "excel", extract_xlsx_text(data))]
    if suffix in {".pptx"}:
        return [fragment(source, "pptx", extract_pptx_text(data))]
    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}:
        return [fragment(source, "image", extract_image_text(data, warnings))]
    if suffix in {".csv"}:
        return [fragment(source, "csv", extract_csv_text(data))]
    if suffix in {".txt", ".md", ".log", ".ini", ".conf", ".json", ".xml", ".html", ".htm", ".eml"}:
        return [fragment(source, "text", decode_text(data))]
    converted = try_libreoffice_text(filename, data)
    if converted:
        return [fragment(source, "converted", converted)]
    return [fragment(source, "binary-text", decode_text(data))]


def fragment(source: SourceFile, kind: str, text: str, location: str = "") -> dict[str, Any]:
    clean = clean_text(text)
    relative_path = (source.relativePath or source.filename or "").replace("\\", "/").strip("/")
    context_text = "\n".join([relative_path, source.filename, clean[:4000]])
    date_hints = extract_date_hints(context_text)
    role = infer_source_role(context_text)
    return {
        "fileId": source.id,
        "filename": source.filename,
        "relativePath": relative_path or source.filename,
        "pathContext": path_context(relative_path or source.filename),
        "sha256": source.sha256,
        "kind": kind,
        "location": location,
        "clientModifiedAt": source.clientModifiedAt,
        "uploadedAt": source.uploadedAt,
        "contentDateHints": date_hints,
        "sourceRole": role,
        "text": clean,
    }


def parse_zip(filename: str, source: SourceFile, data: bytes, warnings: list[str]) -> list[dict[str, Any]]:
    fragments: list[dict[str, Any]] = []
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        for info in archive.infolist():
            name = info.filename.replace("\\", "/")
            if info.is_dir() or name.startswith("/") or ".." in Path(name).parts:
                continue
            if name.lower().endswith(".dmp"):
                warnings.append(f"Rejected dmp file inside {filename}: {name}")
                continue
            if info.file_size > 100 * 1024 * 1024:
                warnings.append(f"Skipped large file inside {filename}: {name}")
                continue
            nested = SourceFile(
                id=source.id,
                filename=f"{filename}/{name}",
                relativePath=f"{source.relativePath or filename}/{name}",
                objectKey=source.objectKey,
                sha256=source.sha256,
                contentType="",
                sizeBytes=info.file_size,
                clientModifiedAt=source.clientModifiedAt,
                uploadedAt=source.uploadedAt,
            )
            try:
                fragments.extend(parse_source(nested.filename, nested, archive.read(info), warnings))
            except Exception as exc:
                warnings.append(f"{nested.filename}: {exc}")
    return fragments


def extract_pdf_text(data: bytes) -> str:
    import fitz

    parts = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for index, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                parts.append(f"[Page {index}]\n{text}")
    return "\n\n".join(parts)


def extract_docx_text(data: bytes) -> str:
    from docx import Document

    document = Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_xlsx_text(data: bytes) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts = []
    for sheet in workbook.worksheets:
        sheet_parts = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if cells:
                sheet_parts.append(" | ".join(cells))
        if sheet_parts:
            parts.append(f"[Sheet: {sheet.title}]")
            parts.extend(sheet_parts)
    return "\n".join(parts)


def extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[Slide {index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def extract_image_text(data: bytes, warnings: list[str]) -> str:
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(image, lang="jpn+chi_sim+eng")
    except Exception as exc:
        warnings.append(f"Image OCR failed: {exc}")
        return ""


def extract_csv_text(data: bytes) -> str:
    text = decode_text(data)
    output = io.StringIO()
    try:
        reader = csv.reader(io.StringIO(text))
        for row in reader:
            if any(cell.strip() for cell in row):
                output.write(" | ".join(cell.strip() for cell in row))
                output.write("\n")
        return output.getvalue()
    except Exception:
        return text


def try_libreoffice_text(filename: str, data: bytes) -> str:
    with tempfile.TemporaryDirectory() as temp_dir:
        source = Path(temp_dir) / filename
        source.parent.mkdir(parents=True, exist_ok=True)
        source.write_bytes(data)
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", temp_dir, str(source)],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=60,
                check=True,
            )
        except Exception:
            return ""
        for candidate in Path(temp_dir).glob("*.txt"):
            if candidate.exists():
                return candidate.read_text(encoding="utf-8", errors="ignore")
    return ""


def decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp932", "shift_jis", "gb18030", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def clean_text(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def path_context(filename: str) -> str:
    normalized = (filename or "").replace("\\", "/").strip("/")
    parts = [part.strip() for part in normalized.split("/") if part.strip()]
    if len(parts) <= 1:
        return parts[0] if parts else ""
    return " > ".join(parts)


def extract_date_hints(text: str) -> list[str]:
    patterns = [
        r"\d{8}\s*(?:以降|以後|以后|以前|まで)?",
        r"\d{4}[-/年]\d{1,2}(?:[-/月]\d{1,2}日?)?\s*(?:以降|以後|以后|以前|まで)?",
        r"(?:以降|以後|以后|以前|まで|旧|新|追加|補足|差分|変更|更新)",
    ]
    hints: list[str] = []
    for pattern in patterns:
        for match in re.finditer(pattern, text, flags=re.IGNORECASE):
            hint = match.group(0).strip()
            if hint and hint not in hints:
                hints.append(hint)
    return hints[:20]


def infer_source_role(text: str) -> str:
    lowered = text.lower()
    if any(token in lowered for token in ["上書き", "置換", "override", "更新", "変更"]):
        return "override"
    if any(token in lowered for token in ["以降", "以後", "以后", "新サーバ", "新vpn", "new"]):
        return "current"
    if any(token in lowered for token in ["追加", "補足", "差分", "supplement"]):
        return "supplement"
    if any(token in lowered for token in ["旧", "以前", "まで", "old", "historical"]):
        return "historical"
    return "unknown"


def parse_datetime(value: str) -> datetime | None:
    if not value:
        return None
    try:
        return datetime.fromisoformat(value.replace("Z", "+00:00"))
    except ValueError:
        return None


def summarize_source_meta(fragments: list[dict[str, Any]]) -> list[dict[str, Any]]:
    by_file: dict[str, dict[str, Any]] = {}
    role_priority = {"current": 5, "override": 4, "supplement": 3, "unknown": 2, "historical": 1}
    for item in fragments:
        file_id = str(item.get("fileId") or item.get("sha256") or item.get("filename") or "")
        meta = by_file.setdefault(file_id, {
            "fileId": item.get("fileId") or "",
            "filename": item.get("filename") or "",
            "relativePath": item.get("relativePath") or item.get("filename") or "",
            "pathContext": item.get("pathContext") or "",
            "clientModifiedAt": item.get("clientModifiedAt") or "",
            "uploadedAt": item.get("uploadedAt") or "",
            "sourceRole": "unknown",
            "dateHints": [],
        })
        role = str(item.get("sourceRole") or "unknown")
        if role_priority.get(role, 0) > role_priority.get(str(meta.get("sourceRole") or "unknown"), 0):
            meta["sourceRole"] = role
        for hint in item.get("contentDateHints") or []:
            if hint not in meta["dateHints"]:
                meta["dateHints"].append(hint)
    return sorted(by_file.values(), key=source_meta_sort_key)


def source_meta_sort_key(item: dict[str, Any]) -> tuple[int, float, str]:
    role_weight = {"current": 0, "override": 1, "supplement": 2, "unknown": 3, "historical": 4}
    modified = parse_datetime(str(item.get("clientModifiedAt") or item.get("uploadedAt") or ""))
    timestamp = modified.timestamp() if modified else 0
    return (role_weight.get(str(item.get("sourceRole") or "unknown"), 3), -timestamp, str(item.get("relativePath") or ""))


def compose_precedence_summary(source_meta: list[dict[str, Any]]) -> str:
    lines = ["Source precedence:"]
    for index, item in enumerate(source_meta, start=1):
        hints = ", ".join(str(hint) for hint in item.get("dateHints") or [])
        lines.append(
            f"{index}. role={item.get('sourceRole') or 'unknown'}; modified={item.get('clientModifiedAt') or '-'}; "
            f"path={item.get('relativePath') or item.get('filename') or '-'}; hints={hints or '-'}"
        )
    return "\n".join(lines)


def compose_raw_text(fragments: list[dict[str, Any]], warnings: list[str]) -> str:
    sections = []
    for item in fragments:
        text = clean_text(str(item.get("text") or ""))
        if not text:
            continue
        path_note = str(item.get("pathContext") or item.get("filename") or "").strip()
        sections.append(
            "\n".join([
                f"===== Source: {item.get('filename') or 'file'} =====",
                f"Path context: {path_note}",
                f"Source role: {item.get('sourceRole') or 'unknown'}",
                f"Client modified: {item.get('clientModifiedAt') or '-'}",
                f"Date hints: {', '.join(str(hint) for hint in item.get('contentDateHints') or []) or '-'}",
                f"Type: {item.get('kind') or 'unknown'}",
                text,
            ])
        )
    if warnings:
        sections.append("===== Hermes Warnings =====\n" + "\n".join(f"- {warning}" for warning in warnings))
    return "\n\n".join(sections).strip()
